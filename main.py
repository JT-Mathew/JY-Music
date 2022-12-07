from pptx import *
from constants import *
from tkinter import Tk
from JYPop import Application
import pandas as pd
import math

#adds verse to slide
def addPara(textBoxText, para):
    splitPara = para.splitlines()
    lineCount = len(splitPara)
    counter = lineCount
    newPara = []
    fontSize = verse_fontSize_lines5
    characterLimit = verse_charLimit_lines5

    if len(splitPara) == 6:
        fontSize = verse_fontSize_lines6
        characterLimit = verse_charLimit_lines6
    elif len(splitPara) > 6:
        fontSize = verse_fontSize_lines7
        characterLimit = verse_charLimit_lines7

    if counter < 6:
        counter = lineCount
        counter = checkLineCount(splitPara, counter, verse_charLimit_lines5)

        if counter >= 6:
            counter = 6

    if counter == 6:
        counter = lineCount
        counter = checkLineCount(splitPara, counter, verse_charLimit_lines6)

        if counter > 6:
            counter = 7
        else:
            fontSize = verse_fontSize_lines6
            characterLimit = verse_charLimit_lines6

    if counter > 6:
        counter = lineCount
        fontSize = verse_fontSize_lines7
        characterLimit = verse_charLimit_lines7

    for line in splitPara:
        line = process_Limit(line, characterLimit)
        newPara.append(line)

    splitPara = newPara

    count = 0

    for line in splitPara:
        addLine(textBoxText, line, count, fontSize)
        count = 1

#adds a line to slide
def addLine(textBoxText, line, count, fontSize):
    if count == 0:
        textBoxPara = textBoxText.paragraphs[0]
    else:
        textBoxPara = textBoxText.add_paragraph()
    textBoxPara.space_after = Pt(10)
    textBoxPara.line_spacing = 0.8

    if fontSize == verse_fontSize_lines5:
        textBoxPara.line_spacing = 0.9
    elif fontSize == verse_fontSize_lines6:
        textBoxPara.line_spacing = 0.8
    else:
        textBoxPara.line_spacing = 0.7

    textBoxPara.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    textBoxPara.font.name = 'Calibri (Body)'
    textBoxPara.font.size = Pt(fontSize)
    textBoxPara.text = line

#checks the line count
def checkLineCount(splitPara, count, limit):
    for line in splitPara:
        if len(line) > limit:
            count = count + 0.7
    count = round(count)
    return count

#adds a jy icon link
def add_hyper_jy(pres, slide, link):
    pic = slide.shapes.add_picture(jy_icon_path, norm_JYIcon_left, 
                                    norm_JYIcon_top, norm_JYIcon_height)
    add_link(pres, pic, link)

#adds a link to a shape
def add_link(pres, shape, link):
    click = shape.click_action
    click.target_slide = pres.slides[link]
    click.action

#adds a textFrame to a slide
def add_Text_Frame(slide, left, top, width, height, vAlign):
    textBox = slide.shapes.add_textbox(left, top, width, height)
    textFrame = textBox.text_frame
    textFrame.vertical_anchor = vAlign

    return textFrame

#adds a heading to a textFrame
def add_Heading(textFrame, fontName, fontSize, lineSpacing, hAlign, text):
    textF_text = textFrame.paragraphs[0]
    textF_text.font.name = fontName
    textF_text.font.size = Pt(fontSize)
    textF_text.line_spacing = lineSpacing
    textF_text.alignment = hAlign

    textF_text.text = text

#adds a shape to a slide
def add_Shape(slide, shapeType, left, top, width, height):
    shape = slide.shapes.add_shape(shapeType, left, top, width, height)
    shape.fill.background()
    shape.line.fill.background()

    return shape

#adds a text to a shape
def add_Shape_Text(shape, vAlign, fontName, fontSize, hAlign):
    textFrame = shape.text_frame
    textFrame.vertical_anchor = vAlign
    indexText = textFrame.paragraphs[0]

    indexText.font.name = fontName
    indexText.font.size = Pt(fontSize)
    indexText.alignment = hAlign

    return indexText

#processes a string through a characterlimit
def process_Limit(line, limit):
    tempName = line
    while len(line) > limit:
        index = line[:limit].rindex(' ')
        tempName = line[:index] + '\n' + line[index:]
        break
    return tempName

#cuts off a title and adds spacing in indexLimit
def index_Limit(line, limit):
    tempName = line
    index1 = line.splitlines()
    if len(index1) >= 2:
        tempName = index1[0]+ ' ' + index1[1]

    while len(tempName) > limit:
        index2 = tempName[:limit].rindex(' ')
        tempName = tempName[:index2]
        break
    return tempName

#removes nan
def clean_Database(df):
    fullDatabase = df.values.tolist()
    x = 0
    for song in fullDatabase:
        fullDatabase[x] = [x for x in song if str(x) != 'nan']
        x = x + 1
    return fullDatabase

#builds the Title Slide
def build_Title_Slide(pres, slide_register, presentationName):
    titleSlide = pres.slides.add_slide(slide_register)
    startBoxText = add_Text_Frame(titleSlide, title_left, title_top, title_width, title_height, title_VAlign)
    add_Heading(startBoxText, title_fontName, title_fontSize, title_lineSpacing, title_HAlign, presentationName)

    subBoxText = add_Text_Frame(titleSlide, subTitle_left, subTitle_top, subTitle_width, subTitle_height, subTitle_VAlign)
    add_Heading(subBoxText, subTitle_fontName, subTitle_fontSize, subTitle_lineSpacing, subTitle_HAlign, subTitle_text)

    pic = titleSlide.shapes.add_picture(jy_icon_path, title_JYIcon_left, title_JYIcon_top, title_JYIcon_height)

#builds the Index Slides
def build_Index_Slide(pres, slide_register, indexSlideIndex):
    indexSlide = pres.slides.add_slide(slide_register)
    indexSlideIndex.append(indexSlide)
    indexTitleTF = add_Text_Frame(indexSlide, indexHeading_left, indexHeading_top, indexHeading_width, indexHeading_height, indexHeading_VAlign)
    add_Heading(indexTitleTF, indexHeading_fontName, indexHeading_fontSize, indexHeading_lineSpacing, indexHeading_HAlign, indexHeading_text)
    add_hyper_jy(pres, indexSlide, 0)

#builds the SongTitle Slides
def build_Song_Title_Slide(pres, slide_register, hyperIndex, presIndex, songName):
    titleSlide = pres.slides.add_slide(slide_register)
    songTitleTF = add_Text_Frame(titleSlide, songTitles_left, songTitles_top, songTitles_width, songTitles_height, songTitles_VAlign)
    add_Heading(songTitleTF, songTitles_fontName, songTitles_fontSize, songTitles_lineSpacing, songTitles_HAlign, songName)
    add_hyper_jy(pres, titleSlide, hyperIndex)
    presIndex.append(pres.slides.index(titleSlide))

#builds the SongVerse Slides
def build_Song_Verse_Slide(pres, slide_register, hyperIndex, verse):
    verseSlide = pres.slides.add_slide(slide_register)
    verseTf = add_Text_Frame(verseSlide, verse_left, verse_top, verse_width, verse_height, verse_VAlign)
    addPara(verseTf, verse)
    add_hyper_jy(pres, verseSlide, hyperIndex)

#populate the Index Slide List
def populate_Index_Slide(pres, indexSlideIndex, indexListLeft, indexListTop, index, chosenSongs):
    indexIndex = 0
    for x in indexSlideIndex:
        for y in indexListLeft:
            for z in indexListTop:
                index_song_shape = add_Shape(x, indexList_Shape, y, z, indexList_width, indexList_height)
                index_song_shape_text = add_Shape_Text(index_song_shape, indexList_VAlign, indexList_fontName, indexList_fontSize, indexList_HAlign)

                try:
                    cutTitle = index_Limit(chosenSongs[indexIndex], indexList_limit)
                    add_link(pres, index_song_shape, index[indexIndex])
                    if indexIndex <9:
                        index_song_shape_text.text = str(indexIndex + 1) + '.    ' + cutTitle
                    elif indexIndex <99:
                        index_song_shape_text.text = str(indexIndex + 1) + '.  ' + cutTitle
                    else:
                        index_song_shape_text.text = str(indexIndex + 1) + '.' + cutTitle
                except:
                    pass

                indexIndex = indexIndex + 1

#main method
def main():
    #pull from database else read from file
    df = pd.read_csv(database_path)
    try:
        url = f'https://docs.google.com/spreadsheets/d/1P3Qu1EQLgcQYWSZQwjY5OWmEnnJMvSSgLkasa6rMC6E/gviz/tq?tqx=out:csv'
        df = pd.read_csv(url)
        df.to_csv(database_path, index=False)
    except:
        df = pd.read_csv(database_path)

    #variables
    allSongs = clean_Database(df)
    fullSongList = df['Song'].tolist()
    indexSlide_index = []
    presIndex = []

    #Create the PopUp Window
    window = Tk()
    window.title("JY Music Slides Generator")
    Application.getWindow(window)
    Application.saveSongList(fullSongList)
    app = Application(master=window)
    app.mainloop()

    #save user inputs: Selected songs, save path, and presentation name for the title slide
    chosenSongs = Application.song_List
    savePath = Application.filepath
    presentationName = process_Limit(Application.presentationName, title_limit)

    #calculate number of Index pages required based on the number of songs chosen
    indexPages = math.ceil(len(chosenSongs)/20)

    #EVERYTHING BELOW BUILDS THE PPT
    #Loads presentation theme
    pr1 = Presentation(presentation_path)
    slide1_register = pr1.slide_layouts[6]

    #Creates Title Slide
    build_Title_Slide(pr1, slide1_register, presentationName)

    #Creates Empty Index Slides
    for x in range(indexPages):
        build_Index_Slide(pr1, slide1_register, indexSlide_index)

    #Creates Song and Verse Slides
    for song in chosenSongs:
        hyperIndex = math.floor(chosenSongs.index(song)/20) + 1
        songName = process_Limit(song, songTitles_limit)
        songIndex = fullSongList.index(song)
        lyrics = allSongs[songIndex][1:]

        build_Song_Title_Slide(pr1, slide1_register, hyperIndex, presIndex, songName)

        for verse in lyrics:
            build_Song_Verse_Slide(pr1, slide1_register, hyperIndex, verse)

    #Updates Index Slides with all the songs added
    populate_Index_Slide(pr1, indexSlide_index, indexList_left, indexList_top, presIndex, chosenSongs)

    if Application.save == 1:
        pr1.save(savePath)

#calls the main method when running the file
if __name__ == "__main__":
    main()