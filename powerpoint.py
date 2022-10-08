from pptx import *
from tkinter import *
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.util import Cm, Pt

#adds verse to slide
def addPara(textBoxText, para):
    splitPara = para.splitlines()
    lineCount = len(splitPara)
    counter = lineCount

    newPara = []
    
    fontSize = Lines5
    characterLimit = Limit5
    if len(splitPara) == 6:
        fontSize = Lines6
        characterLimit = Limit6
    elif len(splitPara) > 6:
        fontSize = Lines7
        characterLimit = Limit7
    
    if counter < 6:
        counter = lineCount
        counter = checkLineCountB(splitPara, counter, Limit5)

        if counter >= 6:
            counter = 6
        
    if counter == 6:
        counter = lineCount
        counter = checkLineCountB(splitPara, counter, Limit6)
        
        if counter > 6:
            counter = 7
        else:
            fontSize = Lines6
            characterLimit = Limit6

    if counter > 6:
        counter = lineCount
        fontSize = Lines7
        characterLimit = Limit7

    newVerse = []
    for line in splitPara:
        while len(line) > characterLimit:
            #lineCount = lineCount + 1
            index = line[:characterLimit].rindex(' ')
            line = line[:index] + '\n' + line[index:]
            break
        newVerse.append(line)
    
    splitPara = newVerse

    count = 0
    for line in splitPara:
        addLine(textBoxText, line, count, fontSize)
        count = 1

#adds a line to slide
def addLine(textBoxText, line, count, fontSize):
    if count == 0:
        textBoxPara = textBoxText.paragraphs[0]
    else:
        textBoxPara = textBoxText.add_paragraph()
    textBoxPara.space_after = Pt(10)
    textBoxPara.line_spacing = 0.8
    textBoxPara.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    textBoxPara.font.name = 'Calibri (Body)'
    textBoxPara.font.size = Pt(fontSize)
    textBoxPara.text = line


def checkLineCountB(splitPara, count, limit):
    check = 0
    for line in splitPara:
        if len(line) > limit:
            if check == 0:
                count = count + 1
                check = 1
            else:
                check = 0
    return count

#to determine font size based on number of lines
Lines5 = 60
Lines6 = 56
Lines7 = 51

#character limit based on font size
Limit5 = 30
Limit6 = 33
Limit7 = 38

#textbox dimensions and location
left = Cm(2.42)
top = Cm(2)
width = Cm(29.01)
height = Cm(14.5)

#verses stored in seperate items in a list
song = []
song.append("""Bless the Lord
Oh my soul, oh my soul
Worship his holy name
Sing like never before
Oh my soul
I’ll worship your holy name""")
song.append("""The sun comes up, it’s a new day dawning
It’s time to sing your song again
Whatever may pass and whatever lies before me
Let me be singing when the evening comes
""")

pr1 = Presentation("MusicSlidesTemplate.pptx")

slide1_register = pr1.slide_layouts[6]



'''
count = 0
slide = pr1.slides.add_slide(slide1_register)
textBox = slide.shapes.add_textbox(left, top, width, height)
textBoxText = textBox.text_frame

for line in newVerse:
    addLine(textBoxText, line, count, fontSize)
    count = 1
'''


for verse in song:
    slide = pr1.slides.add_slide(slide1_register)
    textBox = slide.shapes.add_textbox(left, top, width, height)
    textBoxText = textBox.text_frame
    addPara(textBoxText, verse)

pr1.save('testPower.pptx')

